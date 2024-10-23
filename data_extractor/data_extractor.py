import os, io, csv  # Import necessary libraries
import docx.document  # Import specific module for working with DOCX files
import pdfplumber  # Import pdfplumber to work with PDF files
import docx  # Import python-docx for DOCX files
from pptx import Presentation  # Import python-pptx for PPTX files
from PIL import Image  # Import PIL to handle images

# Universal Data Extractor class to handle different file types (PDF, DOCX, PPTX)
class UniversalDataExtractor():
    def __init__(self, loader):
        """
        Initialize the UniversalDataExtractor with a file loader.
        
        Args:
            loader: An instance of a file loader that handles file loading.
        """
        self.file_loader = loader  # Store the file loader object
        self.content = self.file_loader.load_file()  # Load the file content
        self.file_type = os.path.splitext(loader.file_path)[1].lower()  # Extract the file extension and convert it to lowercase
        
        # Handle different file types (PDF, DOCX, PPTX)
        if self.file_type == '.pdf':
            self.pdf = pdfplumber.open(self.file_loader.file_path)  # Open PDF using pdfplumber
            
        elif self.file_type == '.docx':
            self.doc = docx.Document(self.file_loader.file_path)  # Open DOCX using python-docx
            
        elif self.file_type == '.pptx':
            self.prs = Presentation(self.file_loader.file_path)  # Open PPTX using python-pptx
 
    def extract_text(self):
        """
        Extract text from the file based on its type.
        
        Returns:
            str: Extracted text.
        """
        extracted_data = {'text': '', 'metadata': {}}  # Placeholder for extracted data
        
        # Extract text from a PDF
        if self.file_type == '.pdf':
            return "\n".join(page.extract_text() or "" for page in self.pdf.pages).strip()  # Extract text from all PDF pages
        
        # Extract text from a DOCX
        elif self.file_type == '.docx':
            return "\n".join(para.text for para in self.doc.paragraphs).strip()  # Extract text from DOCX paragraphs
        
        # Extract text from a PPTX
        elif self.file_type == 'pptx':
            text = ''
            for slide in self.prs.slides:  # Iterate through each slide
                for shape in slide.shapes:  # Iterate through each shape
                    if hasattr(shape, "text"):  # Check if the shape has text
                        text += shape.text + "\n"
            return text.strip()  # Return extracted text

        return ""  # Return empty string if file type is not supported
    
    def extract_tables(self):
        """
        Extract tables from the file based on its type.
        
        Returns:
            list: Extracted tables.
        """
        # Extract tables from a PDF
        if self.file_type == ".pdf":
            return [table for page in self.pdf.pages for table in page.extract_tables()]  # Extract tables from all pages of the PDF
        
        # Extract tables from a DOCX
        elif self.file_type == "docx":
            tables = []
            for table in self.doc.tables:  # Iterate through all tables
                table_data = [[cell.text for cell in row.cells] for row in table.rows]  # Extract text from each cell
                tables.append(table_data)
            return tables
        
        # PPTX typically does not contain tables in this context
        elif self.file_type == ".pptx":
            return []
        
        return []
    
    def extract_images(self):
        """
        Extract images from the file based on its type.
        
        Returns:
            list: Paths of extracted images.
        """
        images = []  # List to store image paths
        
        # Extract images from a PDF
        if self.file_type == ".pdf":
            for page_number, page in enumerate(self.pdf.pages):  # Iterate through each page
                for img_index, img in enumerate(page.images):  # Iterate through images in the page
                    if 'stream' in img:  # Check if image has raw data stream
                        img_data = img['stream'].get_rawdata()
                        images.append(self.save_image(img_data, img_index + 1, ".pdf", page_number))  # Save the image and append its path
        
        # Extract images from a DOCX
        elif self.file_type == ".docx":
            for rel in self.doc.part.rels.values():  # Iterate through relationships to find images
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob
                    images.append(self.save_image(img_data, len(images) + 1, ".docx"))  # Save the image and append its path
        
        # Extract images from a PPTX
        elif self.file_type == ".pptx":
            for slide_number, slide in enumerate(self.prs.slides):  # Iterate through each slide
                for shape_index, shape in enumerate(slide.shapes):  # Iterate through shapes
                    if shape.shape_type == 13:  # Shape type 13 is for images
                        img_data = shape.image.blob
                        images.append(self.save_image(img_data, shape_index + 1, ".pptx", slide_number))  # Save the image and append its path
        
        return images
    
    def extract_metadata(self):
        """
        Extract metadata from the file.
        
        Returns:
            dict: Extracted metadata.
        """
        # Extract metadata from a PDF
        if self.file_type == ".pdf":
            return self.pdf.metadata
        
        # Extract metadata from a DOCX
        elif self.file_type == ".docx":
            return self.doc.core_properties
        
        # Extract metadata from a PPTX
        elif self.file_type == ".pptx":
            return self.prs.core_properties
        
        return {}
    
    def extract_links(self):
        """
        Extract hyperlinks from the file.
        
        Returns:
            list: List of extracted links.
        """
        links = []  # List to store extracted links
        
        # Extract links from a PDF
        if self.file_type == ".pdf":
            for page in self.pdf.pages:
                links.extend(annot.get("uri") for annot in getattr(page, 'annots', []) if annot.get("uri"))  # Extract links from annotations
        
        # Extract links from a DOCX
        elif self.file_type == ".docx":
            for rel in self.doc.part.rels.values():
                if "hyperlink" in rel.reltype:
                    links.append(rel.target_ref)  # Extract links from DOCX relationships
        
        # Extract links from a PPTX
        elif self.file_type == ".pptx":
            for slide in self.prs.slides:  # Iterate through slides
                for shape in slide.shapes:  # Iterate through shapes
                    if shape.has_text_frame:  # Check if shape contains text
                        for paragraph in shape.text_frame.paragraphs:  # Iterate through paragraphs
                            for run in paragraph.runs:  # Iterate through runs
                                if run.hyperlink:  # Check if the run has a hyperlink
                                    links.append(run.hyperlink.address)  # Append the hyperlink address
        
        return links
    
    def get_file_name(self):
        """Get the file name from the file path."""
        return os.path.basename(self.file_loader.file_path)
    
    def save_image(self, img_data, index, file_ext, page_number=None):
        """
        Save the extracted image to the output directory and return the image path.
        
        Args:
            img_data (bytes): The raw image data.
            index (int): The image index.
            file_ext (str): The file extension to create unique names.
            page_number (int, optional): The page number for PDFs.
        
        Returns:
            str: The path of the saved image.
        """
        # Create a unique image filename based on the file name, page number (if applicable), and image index
        img_filename = f"{self.get_file_name().replace(file_ext, '')}_img_{index}.png"
        if page_number is not None:
            img_filename = f"{self.get_file_name().replace(file_ext, '')}_page_{page_number + 1}_img_{index}.png"
        
        img_path = os.path.join('output', img_filename)  # Set the output path for the image
        image = Image.open(io.BytesIO(img_data))  # Open the image data using PIL
        image.save(img_path)  # Save the image to the specified path
        return img_path  # Return the saved image path
    
    def close(self):
        """Close the file if it's a PDF."""
        if self.file_type == '.pdf':
            self.pdf.close()  # Close the PDF object
